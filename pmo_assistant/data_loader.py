from __future__ import annotations

from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import Dict, List, Optional

import pandas as pd
from docx import Document
from pptx import Presentation

from .config import BASE_DIR, DOCX_FILES, EXCEL_FILES, PPTX_FILES, resolve_file


def _ensure_exists(path: Path) -> Path:
    if not path.exists():
        raise FileNotFoundError(f"Expected file not found: {path}")
    return path


def _read_excel_sheet_fuzzy(path: Path, required_keywords: List[str]) -> pd.DataFrame:
    """
    Read the first sheet whose normalized name contains all required keywords.
    """
    path = _ensure_exists(path)
    xls = pd.ExcelFile(path)
    keywords = [k.replace(" ", "").lower() for k in required_keywords]

    for sheet_name in xls.sheet_names:
        norm = sheet_name.replace(" ", "").lower()
        if all(k in norm for k in keywords):
            return pd.read_excel(xls, sheet_name=sheet_name)

    raise ValueError(
        f"Could not find sheet in {path.name} with keywords {required_keywords}. "
        f"Available sheets: {xls.sheet_names}"
    )


@lru_cache(maxsize=1)
def load_talent_pool() -> pd.DataFrame:
    """
    Load talent pool data from the Excel workbook.
    """
    excel_path = resolve_file(EXCEL_FILES["talent_pool"])
    # Look for a sheet like "Talent Pool"
    df = _read_excel_sheet_fuzzy(excel_path, ["talent", "pool"])
    return df


@lru_cache(maxsize=1)
def load_resource_allocation_master() -> pd.DataFrame:
    """
    Load overall project resource allocation data.
    """
    excel_path = resolve_file(EXCEL_FILES["resource_allocation"])
    df = _read_excel_sheet_fuzzy(excel_path, ["overall", "project", "res"])
    return df


@lru_cache(maxsize=1)
def load_resource_utilization_master() -> pd.DataFrame:
    """
    Load utilization master data.
    """
    excel_path = resolve_file(EXCEL_FILES["resource_utilization"])
    df = _read_excel_sheet_fuzzy(excel_path, ["master"])
    return df


@lru_cache(maxsize=1)
def load_third_party_resources() -> pd.DataFrame:
    """
    Load third party resources data.
    """
    excel_path = resolve_file(EXCEL_FILES["third_party"])
    # Try to locate a core sheet with "resources" or "3rd" in the name
    try:
        df = _read_excel_sheet_fuzzy(excel_path, ["3rd", "party"])
    except Exception:
        df = _read_excel_sheet_fuzzy(excel_path, ["resource"])
    return df


@lru_cache(maxsize=1)
def load_csat_data() -> pd.DataFrame:
    excel_path = resolve_file(EXCEL_FILES["csat"])
    # Assume main CSAT sheet has "CSAT" in sheet name; fall back to first sheet
    try:
        return _read_excel_sheet_fuzzy(excel_path, ["csat"])
    except Exception:
        xls = pd.ExcelFile(excel_path)
        return pd.read_excel(xls, sheet_name=xls.sheet_names[0])


@lru_cache(maxsize=1)
def load_dashboard_dictionary() -> Dict[str, pd.DataFrame]:
    """
    Load all sheets from the Dashboard Data Dictionary workbook into a dict.
    """
    excel_path = resolve_file(EXCEL_FILES["dashboard_ddit"])
    excel_path = _ensure_exists(excel_path)
    xls = pd.ExcelFile(excel_path)
    return {name: pd.read_excel(xls, sheet_name=name) for name in xls.sheet_names}


@lru_cache(maxsize=1)
def load_rca_data() -> pd.DataFrame:
    excel_path = resolve_file(EXCEL_FILES["rca"])
    df = _read_excel_sheet_fuzzy(excel_path, ["rag"])
    return df


@lru_cache(maxsize=1)
def load_project_plan_sheets() -> Dict[str, pd.DataFrame]:
    """
    Load all sheets from the Project Plan workbook into a dict.
    """
    excel_path = resolve_file(EXCEL_FILES["project_plan"])
    excel_path = _ensure_exists(excel_path)
    xls = pd.ExcelFile(excel_path)
    return {name: pd.read_excel(xls, sheet_name=name) for name in xls.sheet_names}


@lru_cache(maxsize=1)
def load_test_case_sheets() -> Dict[str, pd.DataFrame]:
    """
    Load all sheets from the Test Cases workbook into a dict.
    """
    excel_path = resolve_file(EXCEL_FILES["test_cases"])
    excel_path = _ensure_exists(excel_path)
    xls = pd.ExcelFile(excel_path)
    return {name: pd.read_excel(xls, sheet_name=name) for name in xls.sheet_names}


def _extract_docx_text(path: Path) -> str:
    """
    Extract readable text from a Word document (paragraphs + simple table rows).
    """
    path = _ensure_exists(path)
    doc = Document(path)
    parts: List[str] = []

    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


@lru_cache(maxsize=None)
def load_docx_text(key: str) -> str:
    """
    Load and cache text from one of the configured DOCX_FILES.
    """
    if key not in DOCX_FILES:
        raise KeyError(f"Unknown DOCX key: {key}")
    path = resolve_file(DOCX_FILES[key])
    return _extract_docx_text(path)


def _extract_pptx_text(path: Path) -> str:
    """
    Extract simple slide-wise text from a PowerPoint file.
    """
    path = _ensure_exists(path)
    prs = Presentation(path)
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = (shape.text or "").strip()
                if txt:
                    parts.append(txt)
    return "\n".join(parts)


@lru_cache(maxsize=None)
def load_pptx_text(key: str) -> str:
    """
    Load and cache text from one of the configured PPTX_FILES.
    """
    if key not in PPTX_FILES:
        raise KeyError(f"Unknown PPTX key: {key}")
    path = resolve_file(PPTX_FILES[key])
    return _extract_pptx_text(path)


@dataclass
class EmployeeSummary:
    emp_id: str
    name: str
    role: Optional[str]
    core_skill: Optional[str]
    location: Optional[str]
    bench_pct: Optional[float]
    status: Optional[str]


def build_employee_summaries(limit: Optional[int] = None) -> List[EmployeeSummary]:
    """
    Small helper primarily for debugging / sanity checks in the UI.
    """
    df = load_talent_pool()
    rows = df.to_dict(orient="records")
    summaries: List[EmployeeSummary] = []
    for row in rows[: limit or len(rows)]:
        summaries.append(
            EmployeeSummary(
                emp_id=str(row.get("Employee Code") or row.get("Emp ID") or ""),
                name=str(row.get("Employee Name") or ""),
                role=row.get("Role"),
                core_skill=row.get("Core Skill"),
                location=row.get("Location"),
                bench_pct=float(row.get("Bench %")) if row.get("Bench %") not in (None, "") else None,
                status=row.get("Status"),
            )
        )
    return summaries

