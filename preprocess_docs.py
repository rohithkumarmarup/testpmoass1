import os
import re
from pathlib import Path

import pandas as pd
from docx import Document
from pptx import Presentation

# CHANGE THIS IF YOUR FOLDER MOVES
BASE_DIR = Path(r"C:\Users\rohit\Downloads\Pmoxponnet")
PROCESSED_DIR = BASE_DIR / "processed"
PROCESSED_DIR.mkdir(exist_ok=True)


def safe_name(name: str) -> str:
    """Make a safe filename part from sheet or file names."""
    name = re.sub(r"[^\w\-]+", "_", name.strip())
    return name.strip("_") or "unknown"


def process_excel(path: Path):
    print(f"Processing Excel: {path.name}")
    try:
        xls = pd.ExcelFile(path)
    except Exception as e:
        print(f"  !! Failed to open {path.name}: {e}")
        return

    for sheet_name in xls.sheet_names:
        try:
            df = pd.read_excel(xls, sheet_name=sheet_name)
        except Exception as e:
            print(f"  !! Failed reading sheet '{sheet_name}' in {path.name}: {e}")
            continue

        out_name = f"{safe_name(path.stem)}__{safe_name(sheet_name)}.csv"
        out_path = PROCESSED_DIR / out_name
        try:
            df.to_csv(out_path, index=False)
            print(f"  -> Saved sheet '{sheet_name}' to {out_path.name}")
        except Exception as e:
            print(f"  !! Failed saving CSV for sheet '{sheet_name}': {e}")


def extract_docx_text(path: Path) -> str:
    doc = Document(path)
    parts = []

    # Paragraphs
    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            parts.append(text)

    # Tables (optional but useful)
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                parts.append(" | ".join(row_text))

    return "\n".join(parts)


def process_docx(path: Path):
    print(f"Processing Word: {path.name}")
    try:
        text = extract_docx_text(path)
    except Exception as e:
        print(f"  !! Failed to read {path.name}: {e}")
        return

    out_name = f"{safe_name(path.stem)}.txt"
    out_path = PROCESSED_DIR / out_name
    try:
        out_path.write_text(text, encoding="utf-8")
        print(f"  -> Saved text to {out_path.name}")
    except Exception as e:
        print(f"  !! Failed saving text for {path.name}: {e}")


def extract_pptx_text(path: Path) -> str:
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    parts.append(txt)
    return "\n".join(parts)


def process_pptx(path: Path):
    print(f"Processing PowerPoint: {path.name}")
    try:
        text = extract_pptx_text(path)
    except Exception as e:
        print(f"  !! Failed to read {path.name}: {e}")
        return

    out_name = f"{safe_name(path.stem)}.txt"
    out_path = PROCESSED_DIR / out_name
    try:
        out_path.write_text(text, encoding="utf-8")
        print(f"  -> Saved text to {out_path.name}")
    except Exception as e:
        print(f"  !! Failed saving text for {path.name}: {e}")


def main():
    print(f"Base folder: {BASE_DIR}")
    for entry in BASE_DIR.iterdir():
        if not entry.is_file():
            continue
        # Skip already-processed or temporary files
        if entry.name.startswith("~") or entry.suffix.lower() == ".py":
            continue

        ext = entry.suffix.lower()
        if ext == ".xlsx":
            process_excel(entry)
        elif ext == ".docx":
            process_docx(entry)
        elif ext == ".pptx":
            process_pptx(entry)

    print("Done. Check the 'processed' folder.")


if __name__ == "__main__":
    main()